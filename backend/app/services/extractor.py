"""
extractor.py
============
WHAT THIS FILE DOES
-------------------
It turns ANY supported file into plain text.

    employee_handbook.pdf   ->  "SECTION 1: ANNUAL LEAVE ENTITLEMENT All..."
    q3_budget.xlsx          ->  "Sheet 'Q3 Budget': Department | Budget | ..."
    a_scanned_contract.pdf  ->  (no text layer -> OCR) -> "AGREEMENT made..."

This is stage one of the pipeline. Nothing downstream -- not chunking, not
embedding, not search -- can happen until the text exists.

THE CENTRAL IDEA: ONE PARSER PER FORMAT
---------------------------------------
There is no universal "read any document" library that does a good job. A PDF
stores glyphs at x/y coordinates; a DOCX is a zip of XML; an XLSX is a grid of
typed cells; an email is a MIME tree. Each needs its own tool:

    PDF        -> PyMuPDF        text per page, plus document metadata
    DOCX       -> python-docx     paragraphs AND tables
    PPTX       -> python-pptx     every shape on every slide
    XLSX/XLS   -> openpyxl        every sheet, row by row
    HTML       -> BeautifulSoup   tags stripped, readable text kept
    EML        -> stdlib email    walks the MIME tree for the plain-text body
    CSV/TXT/MD -> stdlib          read directly
    ZIP        -> stdlib zipfile  unpacked, then each file sent back through
                                  this same dispatcher, recursively
    images     -> Tesseract OCR   there is no text to extract, only pixels

`extract()` at the top is the dispatcher that routes a file to the right one.

THE SCANNED-PDF PROBLEM
-----------------------
Some PDFs contain no text at all -- they are photographs of pages. Ask PyMuPDF
for their text and it politely returns nothing.

So after extracting, we MEASURE. A real text page holds 1,500-3,000 characters;
a scan yields close to zero. If a PDF averages fewer than
`OCR_CHAR_THRESHOLD_PER_PAGE` characters per page, we conclude there is no text
layer, rasterise each page into an image, and run OCR instead.

That single heuristic is what lets a scanned contract become just as searchable
as a digital one.
"""

import csv
import os
import uuid
import zipfile
from email import message_from_bytes
from email.policy import default as default_email_policy
from typing import Any, Dict

import fitz  # PyMuPDF. The package installs as `pymupdf` but imports as `fitz`.
import openpyxl
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from langdetect import detect
from pptx import Presentation

from app.core.config import settings
from app.services.ocr import OCRService


class ExtractionResult:
    """
    The uniform result every extractor returns, whatever the input format.

    Having ONE shape means the ingestion task does not need to know or care
    which parser ran -- it just reads `.text`, `.metadata`, `.language`.

    Attributes:
        text      -- the extracted plain text.
        metadata  -- format-specific facts (page count, sheet names, author...).
        language  -- detected two-letter language code, e.g. "en".
        needs_ocr -- True if OCR was used, so the UI can flag it.
    """

    def __init__(
        self,
        text: str,
        metadata: Dict[str, Any],
        language: str,
        needs_ocr: bool = False,
    ):
        self.text = text
        self.metadata = metadata
        self.language = language
        self.needs_ocr = needs_ocr


class DocumentExtractor:
    """Routes each file type to its dedicated parser."""

    # ==================================================================
    # LANGUAGE DETECTION
    # ==================================================================

    @staticmethod
    def detect_language(text: str) -> str:
        """
        Work out which human language the text is written in.

        Stored as metadata. Useful for filtering, for reporting, and as
        groundwork for language-specific OCR or embedding models later.
        """
        # Too little text to judge. langdetect is unreliable on short strings
        # and would essentially guess at random.
        if not text or len(text.strip()) < 10:
            return "en"
        try:
            return detect(text)
        except Exception:
            # langdetect raises on input with no detectable words at all --
            # pure punctuation or numbers. English is a safe default here.
            return "en"

    # ==================================================================
    # THE DISPATCHER
    # ==================================================================

    @classmethod
    def extract(cls, file_path: str, file_type: str) -> ExtractionResult:
        """
        Send the file to whichever parser handles its extension.

        `file_type` is the extension without the dot, e.g. "pdf", already
        lowercased by the upload endpoint -- but we lowercase again here
        because ZIP extraction also calls this method with raw extensions.
        """
        file_type = file_type.lower()

        if file_type == "pdf":
            return cls.extract_pdf(file_path)
        elif file_type in ["png", "jpg", "jpeg", "tiff", "bmp", "gif"]:
            return cls.extract_image(file_path, file_type)
        elif file_type == "docx":
            return cls.extract_docx(file_path)
        elif file_type == "pptx":
            return cls.extract_pptx(file_path)
        elif file_type in ["xlsx", "xls"]:
            return cls.extract_xlsx(file_path)
        elif file_type in ["html", "htm"]:
            return cls.extract_html(file_path)
        elif file_type == "eml":
            return cls.extract_eml(file_path)
        elif file_type in ["txt", "md", "markdown"]:
            return cls.extract_txt_or_md(file_path)
        elif file_type == "csv":
            return cls.extract_csv(file_path)
        elif file_type == "zip":
            return cls.extract_zip(file_path)
        else:
            # Raised rather than returned, so the ingestion task marks the
            # document "failed" with a message the user can actually act on.
            raise ValueError(f"Unsupported file type: {file_type}")

    # ==================================================================
    # PDF  -- the most important format, and the only one needing OCR logic
    # ==================================================================

    @classmethod
    def extract_pdf(cls, file_path: str) -> ExtractionResult:
        """
        Extract text from a PDF, falling back to OCR when it is a scan.

        Pages are joined with an explicit "--- PAGE BREAK ---" marker. That
        marker is not decoration: chunker.py splits on it so every chunk knows
        which page it came from, which is what makes citations say "page 12".
        """
        text = ""
        metadata: Dict[str, Any] = {}
        needs_ocr = False

        try:
            document = fitz.open(file_path)

            # Document-level properties, useful for display and debugging.
            metadata = {
                "title": document.metadata.get("title", ""),
                "author": document.metadata.get("author", ""),
                "creator": document.metadata.get("creator", ""),
                "page_count": len(document),
            }

            pages_text = [page.get_text() for page in document]
            text = "\n--- PAGE BREAK ---\n".join(pages_text)

            # ---- THE SCANNED-DOCUMENT HEURISTIC ----
            # Compare total characters against a per-page budget. A normal text
            # page yields well over the threshold; a scan yields nearly zero.
            page_count = len(document)
            threshold = settings.OCR_CHAR_THRESHOLD_PER_PAGE * page_count

            if page_count > 0 and len(text.strip()) < threshold:
                print(
                    f"[Extractor] Only {len(text.strip())} characters across "
                    f"{page_count} pages -- below the {threshold} threshold. "
                    f"Treating as a scan and running OCR."
                )
                ocr_text, _ = OCRService.ocr_pdf(file_path)
                text = ocr_text
                needs_ocr = True
                metadata["ocr_applied"] = True
                metadata["ocr_provider"] = settings.OCR_PROVIDER

            # Release the file handle. Without this the file stays locked on
            # Windows and cannot be deleted or re-processed.
            document.close()

        except Exception as exc:
            # The PDF is malformed or encrypted in a way PyMuPDF cannot parse.
            # OCR works on the rendered pixels and often succeeds anyway, so it
            # is worth attempting before giving up.
            print(f"[Extractor] PDF parsing failed: {exc}. Attempting OCR...")
            try:
                ocr_text, _ = OCRService.ocr_pdf(file_path)
                text = ocr_text
                needs_ocr = True
                metadata["ocr_applied"] = True
                metadata["ocr_provider"] = settings.OCR_PROVIDER
                metadata["parse_error"] = str(exc)
            except Exception as ocr_error:
                text = f"PDF extraction and OCR fallback both failed: {ocr_error}"
                needs_ocr = True
                metadata["error"] = str(ocr_error)

        return ExtractionResult(
            text=text,
            metadata=metadata,
            language=cls.detect_language(text),
            needs_ocr=needs_ocr,
        )

    # ==================================================================
    # IMAGES  -- pure OCR, there is no text to extract
    # ==================================================================

    @classmethod
    def extract_image(cls, file_path: str, file_type: str) -> ExtractionResult:
        """
        Read the text out of a photograph or screenshot.

        Unlike the PDF path there is no heuristic here: an image never has a
        text layer, so OCR is the only option.
        """
        try:
            # Read as bytes, because OCRService accepts raw bytes regardless of
            # whether they came from a file or from a rasterised PDF page.
            with open(file_path, "rb") as handle:
                image_bytes = handle.read()

            text = OCRService.ocr_image(image_bytes)
            metadata = {
                "image_ocr": True,
                "ocr_provider": settings.OCR_PROVIDER,
                "image_format": file_type,
            }
        except Exception as exc:
            text = f"Image OCR extraction failed: {exc}"
            metadata = {"error": str(exc)}

        return ExtractionResult(
            text=text,
            metadata=metadata,
            language=cls.detect_language(text),
            needs_ocr=True,
        )

    # ==================================================================
    # WORD
    # ==================================================================

    @classmethod
    def extract_docx(cls, file_path: str) -> ExtractionResult:
        """
        Extract a Word document's paragraphs AND its tables.

        Tables matter and are easy to miss: `doc.paragraphs` does NOT include
        text inside table cells. A policy document whose entire retention
        schedule lives in a table would otherwise be extracted as if that
        schedule did not exist.
        """
        text = ""
        metadata: Dict[str, Any] = {}

        try:
            document = DocxDocument(file_path)

            paragraphs = [p.text for p in document.paragraphs]

            # Flatten each table into pipe-separated rows. Simple, but it keeps
            # cells on the same row together, which preserves the association
            # between a label and its value.
            table_rows = []
            for table in document.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))

            text = "\n".join(paragraphs)
            if table_rows:
                text += "\n\nTables:\n" + "\n".join(table_rows)

            metadata = {
                "paragraph_count": len(paragraphs),
                "table_count": len(document.tables),
            }
        except Exception as exc:
            text = f"DOCX extraction failed: {exc}"
            metadata = {"error": str(exc)}

        return ExtractionResult(
            text=text, metadata=metadata, language=cls.detect_language(text)
        )

    # ==================================================================
    # POWERPOINT
    # ==================================================================

    @classmethod
    def extract_pptx(cls, file_path: str) -> ExtractionResult:
        """
        Pull text from every shape on every slide.

        A slide is a canvas of independent shapes, not a flowing document, so
        we walk the shapes and keep only those that carry text. Labelling each
        group "Slide N" gives the model positional context it would otherwise
        lose completely.
        """
        text = ""
        metadata: Dict[str, Any] = {}

        try:
            presentation = Presentation(file_path)

            slides_text = []
            for index, slide in enumerate(presentation.slides, 1):
                shape_texts = []
                for shape in slide.shapes:
                    # Not every shape has text -- pictures, lines and charts do
                    # not -- so check before reading, and skip empty ones.
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_texts.append(shape.text.strip())

                slides_text.append(f"Slide {index}:\n" + "\n".join(shape_texts))

            text = "\n\n".join(slides_text)
            metadata = {"slide_count": len(presentation.slides)}
        except Exception as exc:
            text = f"PPTX extraction failed: {exc}"
            metadata = {"error": str(exc)}

        return ExtractionResult(
            text=text, metadata=metadata, language=cls.detect_language(text)
        )

    # ==================================================================
    # EXCEL
    # ==================================================================

    @classmethod
    def extract_xlsx(cls, file_path: str) -> ExtractionResult:
        """
        Read every sheet of a workbook, row by row.

        `data_only=True` is the key argument: it returns the CALCULATED result
        of a formula rather than the formula text. Without it a total cell
        extracts as "=SUM(B2:B10)", which is useless for answering a question
        about the total.
        """
        text = ""
        metadata: Dict[str, Any] = {}

        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)

            sheets_text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                rows_text = []
                # `values_only=True` yields plain tuples of cell values instead
                # of Cell objects, which is both faster and simpler here.
                for row in sheet.iter_rows(values_only=True):
                    # `any(row)` skips entirely blank rows, which spreadsheets
                    # are full of and which add nothing but noise.
                    if any(row):
                        cells = [
                            str(value) if value is not None else "" for value in row
                        ]
                        rows_text.append(" | ".join(cells))

                sheets_text.append(f"Sheet '{sheet_name}':\n" + "\n".join(rows_text))

            text = "\n\n".join(sheets_text)
            metadata = {"sheet_names": workbook.sheetnames}
        except Exception as exc:
            text = f"Excel extraction failed: {exc}"
            metadata = {"error": str(exc)}

        return ExtractionResult(
            text=text, metadata=metadata, language=cls.detect_language(text)
        )

    # ==================================================================
    # HTML
    # ==================================================================

    @classmethod
    def extract_html(cls, file_path: str) -> ExtractionResult:
        """
        Strip the tags and keep the readable text.

        `get_text(separator="\\n")` inserts a newline between elements, which
        stops adjacent block elements running together into one unreadable
        line -- "Contact usAbout usPricing" rather than three separate items.
        """
        text = ""
        metadata: Dict[str, Any] = {}

        try:
            # `errors="ignore"` because scraped HTML frequently contains bytes
            # that do not decode as UTF-8. Losing a stray character is far
            # better than failing the whole document.
            with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
                soup = BeautifulSoup(handle.read(), "html.parser")

            raw_text = soup.get_text(separator="\n")
            # Drop blank lines and trim each remaining one, since HTML is full
            # of indentation whitespace that carries no meaning.
            text = "\n".join(
                line.strip() for line in raw_text.splitlines() if line.strip()
            )
            metadata = {"title": soup.title.string if soup.title else ""}
        except Exception as exc:
            text = f"HTML extraction failed: {exc}"
            metadata = {"error": str(exc)}

        return ExtractionResult(
            text=text, metadata=metadata, language=cls.detect_language(text)
        )

    # ==================================================================
    # EMAIL
    # ==================================================================

    @classmethod
    def extract_eml(cls, file_path: str) -> ExtractionResult:
        """
        Extract an email's headers and its plain-text body.

        An email is a TREE, not a flat file: a typical message carries both a
        text/plain and a text/html version of the same content, plus any
        attachments. We walk that tree and take only the plain-text parts,
        skipping attachments -- which would otherwise dump base64 gibberish
        straight into the search index.
        """
        text = ""
        metadata: Dict[str, Any] = {}

        try:
            with open(file_path, "rb") as handle:
                message = message_from_bytes(
                    handle.read(), policy=default_email_policy
                )

            body = ""
            if message.is_multipart():
                # `.walk()` visits every part of the MIME tree recursively.
                for part in message.walk():
                    content_type = part.get_content_type()
                    disposition = str(part.get_params())

                    if content_type == "text/plain" and "attachment" not in disposition:
                        body += part.get_payload(decode=True).decode(errors="ignore")
            else:
                # A simple single-part message.
                body = message.get_payload(decode=True).decode(errors="ignore")

            metadata = {
                "subject": message.get("Subject", ""),
                "from": message.get("From", ""),
                "to": message.get("To", ""),
                "date": message.get("Date", ""),
            }

            # Headers are prepended to the body so that a question like "who
            # sent this?" can be answered from the chunk text itself.
            text = (
                f"From: {metadata['from']}\n"
                f"To: {metadata['to']}\n"
                f"Date: {metadata['date']}\n"
                f"Subject: {metadata['subject']}\n\n"
                f"Body:\n{body}"
            )
        except Exception as exc:
            text = f"Email extraction failed: {exc}"
            metadata = {"error": str(exc)}

        return ExtractionResult(
            text=text, metadata=metadata, language=cls.detect_language(text)
        )

    # ==================================================================
    # PLAIN TEXT AND MARKDOWN
    # ==================================================================

    @classmethod
    def extract_txt_or_md(cls, file_path: str) -> ExtractionResult:
        """
        Read the file as-is. Markdown is deliberately left unrendered so the
        `#` headers survive for the markdown chunking strategy to split on.
        """
        text = ""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
                text = handle.read()
        except Exception as exc:
            text = f"Text extraction failed: {exc}"

        return ExtractionResult(
            text=text, metadata={}, language=cls.detect_language(text)
        )

    # ==================================================================
    # CSV
    # ==================================================================

    @classmethod
    def extract_csv(cls, file_path: str) -> ExtractionResult:
        """
        Flatten a CSV into readable lines, one per row.

        We use the `csv` module rather than splitting on commas, because it
        correctly handles quoted fields that themselves contain commas -- a
        naive split would mangle every such row.
        """
        text = ""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
                reader = csv.reader(handle)
                rows = [" , ".join(row) for row in reader]
                text = "\n".join(rows)
        except Exception as exc:
            text = f"CSV extraction failed: {exc}"

        return ExtractionResult(
            text=text, metadata={}, language=cls.detect_language(text)
        )

    # ==================================================================
    # ZIP  -- the recursive one
    # ==================================================================

    @classmethod
    def extract_zip(cls, file_path: str) -> ExtractionResult:
        """
        Unpack an archive and extract every supported file inside it.

        This is the only RECURSIVE extractor: for each file it finds, it calls
        `cls.extract(...)` again, which routes that file to its own parser. So
        a zip containing a PDF, a spreadsheet and a Word document is handled
        correctly with no special-casing at all.

        Uploading one archive of a whole folder is far more practical than
        uploading forty files individually.
        """
        texts = []
        metadata: Dict[str, Any] = {"files_extracted": []}
        temp_dir = None

        try:
            # A uuid in the directory name so two archives extracted at the same
            # moment by two workers cannot overwrite each other's files.
            temp_dir = os.path.join(
                os.path.dirname(file_path), f"zip_temp_{uuid.uuid4().hex}"
            )
            os.makedirs(temp_dir, exist_ok=True)

            with zipfile.ZipFile(file_path, "r") as archive:
                archive.extractall(temp_dir)

            # Formats we know how to read. Anything else inside the archive --
            # executables, videos, fonts -- is skipped silently.
            supported = {
                "pdf", "docx", "pptx", "xlsx", "xls", "html", "htm",
                "eml", "txt", "md", "markdown", "csv",
                "png", "jpg", "jpeg", "tiff", "bmp", "gif",
            }

            # `os.walk` descends into nested folders too, so a zip containing
            # directories is handled correctly.
            for root, _, filenames in os.walk(temp_dir):
                for filename in filenames:
                    inner_path = os.path.join(root, filename)
                    extension = os.path.splitext(filename)[1].lstrip(".").lower()

                    if extension not in supported:
                        continue

                    try:
                        # THE RECURSIVE CALL. Each inner file goes back through
                        # the dispatcher and is handled by its own parser.
                        result = cls.extract(inner_path, extension)
                        texts.append(f"--- File: {filename} ---\n{result.text}")
                        metadata["files_extracted"].append(filename)
                    except Exception as inner_error:
                        # One unreadable file must not fail the whole archive.
                        texts.append(
                            f"--- File: {filename} (extraction failed: {inner_error}) ---"
                        )

            text = "\n\n=================================\n\n".join(texts)
            metadata["file_count"] = len(metadata["files_extracted"])

        except Exception as exc:
            text = f"ZIP extraction failed: {exc}"
            metadata["error"] = str(exc)

        finally:
            # `finally` runs whether we succeeded or failed, so extracted files
            # are always cleaned up and never accumulate on disk.
            if temp_dir and os.path.isdir(temp_dir):
                cls._remove_tree(temp_dir)

        return ExtractionResult(
            text=text, metadata=metadata, language=cls.detect_language(text)
        )

    @staticmethod
    def _remove_tree(directory: str) -> None:
        """
        Delete a directory and everything inside it.

        `topdown=False` makes os.walk yield the DEEPEST directories first,
        which is required: a directory can only be removed once it is empty,
        so children must go before their parents.
        """
        try:
            for root, dirs, files in os.walk(directory, topdown=False):
                for name in files:
                    try:
                        os.remove(os.path.join(root, name))
                    except OSError:
                        pass
                for name in dirs:
                    try:
                        os.rmdir(os.path.join(root, name))
                    except OSError:
                        pass
            os.rmdir(directory)
        except OSError:
            # Cleanup is best-effort. A leftover temp folder is untidy but is
            # not worth failing an otherwise successful extraction over.
            pass
